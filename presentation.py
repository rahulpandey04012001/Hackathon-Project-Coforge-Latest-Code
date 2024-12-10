from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create a presentation object
prs = Presentation()


# Function to add a slide with title and content
def add_slide(prs, title, content, bullet_points=None):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    # Set the title
    title_placeholder = slide.shapes.title
    title_placeholder.text = title

    # Set the content
    content_placeholder = slide.placeholders[1]
    text_frame = content_placeholder.text_frame
    text_frame.clear()
    p = text_frame.add_paragraph()
    p.text = content
    p.font.size = Pt(18)

    # Add bullet points if any
    if bullet_points:
        for point in bullet_points:
            p = text_frame.add_paragraph()
            p.text = point
            p.level = 1


# Title Slide
slide_layout = prs.slide_layouts[0]
title_slide = prs.slides.add_slide(slide_layout)
title = title_slide.shapes.title
title.text = "Requirement Rebels - Automated Test Case Generator"
subtitle = title_slide.placeholders[1]
subtitle.text = "Challenging Limits, Defining Quality\nTeam: Requirement Rebels"

# Slide 2: The Problem
add_slide(
    prs,
    "The Problem We Are Solving",
    "Current Scenario: Manual test case generation delays the software cycle, leading to inconsistent coverage and increased costs.",
    [
        "Manual test creation is time-consuming.",
        "Negative scenarios are frequently missed.",
        "Slow software release cycles."
    ]
)

# Slide 3: How Our Solution Helps
add_slide(
    prs,
    "How Our Solution Helps",
    "Our solution automates test case generation, saving time and increasing consistency.",
    [
        "Automated test generation reduces time requirements.",
        "Focus on exploratory testing improves QA productivity.",
        "Better coverage for industries like software, finance, healthcare, and retail."
    ]
)

# Slide 4: Stakeholders Who Benefit
add_slide(
    prs,
    "Stakeholders Who Benefit",
    "Our solution impacts a range of stakeholders, enhancing their efficiency and reducing costs.",
    [
        "Quality Assurance Teams: More time for complex scenarios.",
        "Product Owners & BAs: Better alignment between user stories and test cases.",
        "Development Teams: Fewer regressions and reliable releases.",
        "Project Managers: Reduced project duration and costs."
    ]
)

# Slide 5: Solution Overview
add_slide(
    prs,
    "Our Solution",
    "An overview of how our solution works to automate the generation of test cases.",
    [
        "Step 1: Input User Story.",
        "Step 2: NLP Parsing to understand requirements.",
        "Step 3: Automated test case generation using LLM.",
        "Step 4: Gherkin feature file generation for BDD frameworks."
    ]
)

# Slide 6: MVP Use Cases
add_slide(
    prs,
    "MVP Use Cases",
    "The key use cases being implemented in the MVP version of our solution.",
    [
        "Single User Story Input: Generates positive and negative test cases.",
        "Multiple User Stories: Upload template to generate test cases for all.",
        "Basic vs. Advanced Comparison: Highlights the impact of LLM-generated cases."
    ]
)

# Slide 7: Future Enhancements
add_slide(
    prs,
    "Future Enhancements",
    "How we plan to enhance the solution in the future.",
    [
        "Integration with CI/CD Pipelines.",
        "Customizable Templates for different industries.",
        "OpenAI Fine-Tuning for precise test cases.",
        "Defect Prediction using AI for better quality assurance."
    ]
)

# Slide 8: Demo Slide
add_slide(
    prs,
    "Live Demo",
    "Watch the power of automated testing in action."
)

# Slide 9: Summary and Closure
add_slide(
    prs,
    "Summary",
    "We addressed the inefficiencies of manual test creation and provided an automated approach to improve software quality.",
    [
        "The Problem: Manual testing is inefficient.",
        "The Solution: Automated test generation using GPT.",
        "The Impact: Improved efficiency, reduced costs, and faster releases."
    ]
)

# Slide 10: Q&A
add_slide(
    prs,
    "Questions & Answers",
    "We believe Requirement Rebels is the right step towards automated, intelligent QA processes. We’re open to any questions."
)

# Save the presentation
prs.save('Requirement_Rebels_Presentation.pptx')

print("Presentation generated successfully!")
